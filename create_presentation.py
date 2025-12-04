from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(46, 134, 171)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(46, 134, 171)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.level = 0
        p.space_before = Pt(6)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Add a slide with an image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(46, 134, 171)
    
    # Add image if it exists
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.2), width=Inches(8))
    
    # Add caption if provided
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_p = caption_frame.paragraphs[0]
        caption_p.font.size = Pt(12)
        caption_p.font.italic = True
        caption_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_code_slide(prs, title, code_text):
    """Add a slide with code"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(46, 134, 171)
    
    # Add code box
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
    text_frame = code_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = code_text
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Courier New'
        paragraph.font.size = Pt(12)
    
    return slide

def add_table_slide(prs, title, data, headers):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(46, 134, 171)
    
    # Add table
    rows = len(data) + 1
    cols = len(headers)
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5)).table
    
    # Set headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(46, 134, 171)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Fill data
    for i, row in enumerate(data):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "IPO Analysis: Questions 5 & 6", 
                "SPAC vs Non-SPAC Returns and Index Inclusion Performance")

# Slide 2: Overview
add_content_slide(prs, "Analysis Overview", [
    "Question 5: Examining SPAC vs Non-SPAC IPO Returns",
    "  • Day 0 return analysis with abnormal return flagging",
    "  • Multi-window return comparison (5-day, 22-day, 91-day, 252-day)",
    "",
    "Question 6: Index Inclusion Performance Analysis",
    "  • S&P 500 inclusion impact on 1-year returns",
    "  • Russell 1000 inclusion impact on 1-year returns",
    "",
    "Dataset: 3,681 IPOs (251 SPACs, 3,430 Non-SPACs)"
])

# Slide 3: Question 5(i) - Code
code_q5i = """# Flag abnormal returns (>= 100%)
stock_ipos['day0_lvl'] = np.where(
    stock_ipos['sym_day0_OTC'] < 1, 
    'normal', 
    'abnormal'
)

# Summary statistics by level and SPAC status
day0_summary = stock_ipos.groupby(
    ['day0_lvl', 'spac']
)['sym_day0_OTC'].agg([
    'mean', 'median', 'std', 'count', 'min', 'max'
])"""

add_code_slide(prs, "Question 5(i): Day 0 Return Analysis - Code", code_q5i)

# Slide 4: Question 5(i) - Results Table
q5i_data = [
    ["Abnormal, Non-SPAC", "2.918", "2.245", "3.010", "30"],
    ["Normal, Non-SPAC", "-0.005", "0.000", "0.127", "3,400"],
    ["Normal, SPAC", "-0.009", "0.000", "0.070", "251"]
]
add_table_slide(prs, "Question 5(i): Day 0 Return Statistics", 
                q5i_data, 
                ["Category", "Mean", "Median", "Std Dev", "Count"])

# Slide 5: Question 5(i) - Visualization
add_image_slide(prs, "Question 5(i): Day 0 Returns - Visual Comparison", 
                "day0_comparison.png",
                "SPACs show lower but more stable Day 0 returns compared to Non-SPACs")

# Slide 6: Question 5(i) - Key Findings
add_content_slide(prs, "Question 5(i): Key Findings", [
    "SPACs have lower and negative mean Day 0 returns:",
    "  • SPACs: -0.88% vs Non-SPACs: +2.09%",
    "",
    "SPACs show significantly lower volatility:",
    "  • SPAC std dev: 0.070 vs Non-SPAC: 0.408",
    "",
    "Abnormal returns (≥100%) only occur in Non-SPACs:",
    "  • 30 cases with mean return of 292%",
    "",
    "Both groups have median returns of 0%:",
    "  • Suggests many IPOs trade at offer price on Day 0"
])

# Slide 7: Question 5(ii) - Code
code_q5ii = """# Analyze returns across multiple windows
windows = ['sym_5day_ret', 'sym_22day_ret', 
           'sym_91day_ret', 'sym_252day_ret']

for window in windows:
    summary = stock_ipos.groupby('spac')[window].agg([
        'mean', 'median', 'std', 'count'
    ])
    print(f"{window} Statistics:")
    print(summary)"""

add_code_slide(prs, "Question 5(ii): Multi-Window Analysis - Code", code_q5ii)

# Slide 8: Question 5(ii) - Results Table
q5ii_data = [
    ["5-day", "0.034", "0.021", "1.252", "0.276"],
    ["22-day", "0.047", "0.030", "1.344", "0.377"],
    ["91-day", "0.061", "0.030", "2.007", "0.477"],
    ["252-day", "0.043", "0.013", "2.968", "0.157"]
]
add_table_slide(prs, "Question 5(ii): Mean Returns & Volatility by Window", 
                q5ii_data,
                ["Window", "Non-SPAC Mean", "SPAC Mean", "Non-SPAC Std", "SPAC Std"])

# Slide 9: Question 5(ii) - Mean Returns Chart
add_image_slide(prs, "Question 5(ii): Mean Returns Across Time Windows", 
                "multiwindow_comparison.png",
                "Non-SPACs consistently outperform SPACs across all time horizons")

# Slide 10: Question 5(ii) - Volatility Chart
add_image_slide(prs, "Question 5(ii): Return Volatility Comparison", 
                "volatility_comparison.png",
                "SPACs demonstrate consistently lower volatility than Non-SPACs")

# Slide 11: Question 6 - Code
code_q6 = """# S&P 500 inclusion performance
sp_performance = stock_ipos.groupby('sp')[
    'sym_252day_ret'
].agg(['mean', 'median', 'std', 'count'])

# Russell 1000 inclusion performance
russell_performance = stock_ipos.groupby('russell')[
    'sym_252day_ret'
].agg(['mean', 'median', 'std', 'count'])"""

add_code_slide(prs, "Question 6: Index Inclusion Analysis - Code", code_q6)

# Slide 12: Question 6 - Results & Visualizations
q6_data = [
    ["Not in S&P 500", "0.037", "0.000", "3,630"],
    ["In S&P 500", "0.293", "0.171", "51"],
    ["Not in Russell 1000", "0.031", "-0.001", "3,538"],
    ["In Russell 1000", "0.282", "0.134", "143"]
]
add_table_slide(prs, "Question 6: Index Inclusion Performance (1-Year Returns)", 
                q6_data,
                ["Category", "Mean Return", "Median Return", "Count"])

# Slide 13: Conclusion
add_content_slide(prs, "Key Conclusions", [
    "Question 5 - SPAC vs Non-SPAC Performance:",
    "  • SPACs underperform across all time horizons",
    "  • SPACs offer lower volatility and more stable returns",
    "  • Non-SPACs have extreme outliers (both positive and negative)",
    "",
    "Question 6 - Index Inclusion Impact:",
    "  • Index inclusion strongly predicts superior performance",
    "  • S&P 500 included: 29.3% vs 3.7% mean 1-year return",
    "  • Russell 1000 included: 28.2% vs 3.1% mean 1-year return",
    "  • Included stocks show lower volatility despite higher returns",
    "  • Very selective: only 1.4% achieve S&P 500, 3.9% Russell 1000"
])

# Save presentation
prs.save('IPO_Analysis_Q5_Q6.pptx')
print("PowerPoint presentation created successfully: IPO_Analysis_Q5_Q6.pptx")
print(f"Total slides: {len(prs.slides)}")
