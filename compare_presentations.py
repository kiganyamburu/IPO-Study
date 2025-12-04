from pptx import Presentation
import os

def analyze_presentation(pptx_path):
    """Analyze a PowerPoint presentation and return detailed information"""
    if not os.path.exists(pptx_path):
        return f"File not found: {pptx_path}"
    
    prs = Presentation(pptx_path)
    
    info = {
        'filename': os.path.basename(pptx_path),
        'total_slides': len(prs.slides),
        'slide_details': []
    }
    
    for i, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': i,
            'layout_name': slide.slide_layout.name,
            'shapes_count': len(slide.shapes),
            'text_content': [],
            'has_images': False,
            'has_tables': False,
            'has_charts': False
        }
        
        for shape in slide.shapes:
            # Check for text
            if hasattr(shape, "text") and shape.text.strip():
                slide_info['text_content'].append(shape.text.strip()[:100])  # First 100 chars
            
            # Check for images
            if shape.shape_type == 13:  # Picture
                slide_info['has_images'] = True
            
            # Check for tables
            if shape.shape_type == 19:  # Table
                slide_info['has_tables'] = True
            
            # Check for charts
            if shape.shape_type == 3:  # Chart
                slide_info['has_charts'] = True
        
        info['slide_details'].append(slide_info)
    
    return info

# Analyze both presentations
print("=" * 80)
print("POWERPOINT COMPARISON ANALYSIS")
print("=" * 80)

file1 = "IPO_Analysis_Q5_Q6.pptx"
file2 = "Your paragraph text.pptx"

print(f"\nAnalyzing: {file1}")
print("-" * 80)
info1 = analyze_presentation(file1)

if isinstance(info1, dict):
    print(f"Total Slides: {info1['total_slides']}")
    print(f"\nSlide-by-slide breakdown:")
    for slide in info1['slide_details']:
        print(f"\n  Slide {slide['slide_number']}:")
        print(f"    Layout: {slide['layout_name']}")
        print(f"    Shapes: {slide['shapes_count']}")
        print(f"    Images: {'Yes' if slide['has_images'] else 'No'}")
        print(f"    Tables: {'Yes' if slide['has_tables'] else 'No'}")
        print(f"    Charts: {'Yes' if slide['has_charts'] else 'No'}")
        if slide['text_content']:
            print(f"    First text: {slide['text_content'][0][:60]}...")
else:
    print(info1)

print("\n" + "=" * 80)
print(f"\nAnalyzing: {file2}")
print("-" * 80)
info2 = analyze_presentation(file2)

if isinstance(info2, dict):
    print(f"Total Slides: {info2['total_slides']}")
    print(f"\nSlide-by-slide breakdown:")
    for slide in info2['slide_details']:
        print(f"\n  Slide {slide['slide_number']}:")
        print(f"    Layout: {slide['layout_name']}")
        print(f"    Shapes: {slide['shapes_count']}")
        print(f"    Images: {'Yes' if slide['has_images'] else 'No'}")
        print(f"    Tables: {'Yes' if slide['has_tables'] else 'No'}")
        print(f"    Charts: {'Yes' if slide['has_charts'] else 'No'}")
        if slide['text_content']:
            print(f"    First text: {slide['text_content'][0][:60]}...")
else:
    print(info2)

# Comparison summary
print("\n" + "=" * 80)
print("COMPARISON SUMMARY")
print("=" * 80)

if isinstance(info1, dict) and isinstance(info2, dict):
    print(f"\nSlide Count:")
    print(f"  {file1}: {info1['total_slides']} slides")
    print(f"  {file2}: {info2['total_slides']} slides")
    print(f"  Difference: {abs(info1['total_slides'] - info2['total_slides'])} slides")
    
    print(f"\nVisual Elements:")
    images1 = sum(1 for s in info1['slide_details'] if s['has_images'])
    images2 = sum(1 for s in info2['slide_details'] if s['has_images'])
    print(f"  {file1}: {images1} slides with images")
    print(f"  {file2}: {images2} slides with images")
    
    tables1 = sum(1 for s in info1['slide_details'] if s['has_tables'])
    tables2 = sum(1 for s in info2['slide_details'] if s['has_tables'])
    print(f"  {file1}: {tables1} slides with tables")
    print(f"  {file2}: {tables2} slides with tables")
    
    print(f"\nAverage shapes per slide:")
    avg_shapes1 = sum(s['shapes_count'] for s in info1['slide_details']) / len(info1['slide_details'])
    avg_shapes2 = sum(s['shapes_count'] for s in info2['slide_details']) / len(info2['slide_details'])
    print(f"  {file1}: {avg_shapes1:.1f}")
    print(f"  {file2}: {avg_shapes2:.1f}")
    
    print("\n" + "=" * 80)
    print("KEY DIFFERENCES")
    print("=" * 80)
    
    if info1['total_slides'] > info2['total_slides']:
        print(f"\n- {file1} has MORE slides ({info1['total_slides']} vs {info2['total_slides']})")
    elif info1['total_slides'] < info2['total_slides']:
        print(f"\n- {file2} has MORE slides ({info2['total_slides']} vs {info1['total_slides']})")
    else:
        print(f"\n- Both presentations have the SAME number of slides ({info1['total_slides']})")
    
    if images1 > images2:
        print(f"- {file1} has MORE images ({images1} vs {images2})")
    elif images1 < images2:
        print(f"- {file2} has MORE images ({images2} vs {images1})")
    
    if tables1 > tables2:
        print(f"- {file1} has MORE tables ({tables1} vs {tables2})")
    elif tables1 < tables2:
        print(f"- {file2} has MORE tables ({tables2} vs {tables1})")

print("\n" + "=" * 80)

# Save to file
with open('comparison_report.txt', 'w', encoding='utf-8') as f:
    f.write("POWERPOINT COMPARISON COMPLETE\n")
    f.write(f"{file1}: {info1['total_slides'] if isinstance(info1, dict) else 'N/A'} slides\n")
    f.write(f"{file2}: {info2['total_slides'] if isinstance(info2, dict) else 'N/A'} slides\n")

print("Report saved to: comparison_report.txt")
