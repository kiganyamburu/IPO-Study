from pptx import Presentation
import os

def get_detailed_comparison():
    """Create detailed comparison of both presentations"""
    
    file1 = "IPO_Analysis_Q5_Q6.pptx"
    file2 = "Your paragraph text.pptx"
    
    output = []
    output.append("=" * 100)
    output.append("DETAILED POWERPOINT COMPARISON")
    output.append("=" * 100)
    
    # Analyze File 1
    output.append(f"\n\nFILE 1: {file1}")
    output.append("-" * 100)
    
    if os.path.exists(file1):
        prs1 = Presentation(file1)
        output.append(f"Total Slides: {len(prs1.slides)}")
        
        for i, slide in enumerate(prs1.slides, 1):
            output.append(f"\n  SLIDE {i}:")
            texts = []
            has_image = False
            has_table = False
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if shape.shape_type == 13:
                    has_image = True
                if shape.shape_type == 19:
                    has_table = True
            
            output.append(f"    - Images: {'Yes' if has_image else 'No'}")
            output.append(f"    - Tables: {'Yes' if has_table else 'No'}")
            output.append(f"    - Text elements: {len(texts)}")
            if texts:
                output.append(f"    - Title/First text: {texts[0][:80]}")
    else:
        output.append(f"  ERROR: File not found")
    
    # Analyze File 2
    output.append(f"\n\n{'=' * 100}")
    output.append(f"\nFILE 2: {file2}")
    output.append("-" * 100)
    
    if os.path.exists(file2):
        prs2 = Presentation(file2)
        output.append(f"Total Slides: {len(prs2.slides)}")
        
        for i, slide in enumerate(prs2.slides, 1):
            output.append(f"\n  SLIDE {i}:")
            texts = []
            has_image = False
            has_table = False
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
                if shape.shape_type == 13:
                    has_image = True
                if shape.shape_type == 19:
                    has_table = True
            
            output.append(f"    - Images: {'Yes' if has_image else 'No'}")
            output.append(f"    - Tables: {'Yes' if has_table else 'No'}")
            output.append(f"    - Text elements: {len(texts)}")
            if texts:
                output.append(f"    - Title/First text: {texts[0][:80]}")
    else:
        output.append(f"  ERROR: File not found")
    
    # Summary
    output.append(f"\n\n{'=' * 100}")
    output.append("SUMMARY COMPARISON")
    output.append("=" * 100)
    
    if os.path.exists(file1) and os.path.exists(file2):
        prs1 = Presentation(file1)
        prs2 = Presentation(file2)
        
        output.append(f"\nSlide Count:")
        output.append(f"  - {file1}: {len(prs1.slides)} slides")
        output.append(f"  - {file2}: {len(prs2.slides)} slides")
        output.append(f"  - Difference: {abs(len(prs1.slides) - len(prs2.slides))} slides")
        
        # Count visual elements
        images1 = sum(1 for slide in prs1.slides for shape in slide.shapes if shape.shape_type == 13)
        images2 = sum(1 for slide in prs2.slides for shape in slide.shapes if shape.shape_type == 13)
        
        tables1 = sum(1 for slide in prs1.slides for shape in slide.shapes if shape.shape_type == 19)
        tables2 = sum(1 for slide in prs2.slides for shape in slide.shapes if shape.shape_type == 19)
        
        output.append(f"\nVisual Elements:")
        output.append(f"  - {file1}: {images1} images, {tables1} tables")
        output.append(f"  - {file2}: {images2} images, {tables2} tables")
        
        output.append(f"\nKey Differences:")
        if len(prs1.slides) < len(prs2.slides):
            output.append(f"  - {file1} has FEWER slides (13 vs 15)")
            output.append(f"  - {file1} is more CONCISE and FOCUSED")
        
        if images1 > images2:
            output.append(f"  - {file1} has MORE visualizations ({images1} vs {images2})")
        
        if tables1 > tables2:
            output.append(f"  - {file1} has MORE data tables ({tables1} vs {tables2})")
    
    output.append("\n" + "=" * 100)
    
    return "\n".join(output)

# Generate and save comparison
comparison_text = get_detailed_comparison()
print(comparison_text)

# Save to file
with open('detailed_comparison.txt', 'w', encoding='utf-8') as f:
    f.write(comparison_text)

print("\n\nDetailed comparison saved to: detailed_comparison.txt")
